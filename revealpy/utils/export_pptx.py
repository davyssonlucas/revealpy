# utils/export_pptx.py
from pptx import Presentation as PPTXPresentation
from pptx.util import Inches
import requests
from io import BytesIO
from PIL import Image
from cairosvg import svg2png

from revealpy.core.content import ContentType


class PPTXExporter:
    def __init__(self, presentation):
        self.presentation = presentation  # Store the presentation object

    def convert_to_pptx(self, output_file="presentation.pptx"):
        pptx = PPTXPresentation()

        for slide in self.presentation.slides:
            ppt_slide = pptx.slides.add_slide(pptx.slide_layouts[1])  # Using layout with title and content

            # Add title
            if ppt_slide.shapes.title:
                ppt_slide.shapes.title.text = slide.title

            # Process each content element
            current_top = Inches(2)  # Start below title
            for content in slide.contents:
                current_top = self._add_content_to_slide(ppt_slide, content, current_top)

        pptx.save(output_file)
        print(f"Arquivo PPTX salvo como {output_file}")

    def _add_content_to_slide(self, ppt_slide, content, top):
        """Add content to PowerPoint slide and return the new top position."""
        if content.type == ContentType.TEXT:
            text_box = ppt_slide.shapes.add_textbox(
                Inches(1), top, Inches(8), Inches(1)
            )
            text_frame = text_box.text_frame
            text_frame.text = content.value
            return top + Inches(1.2)

        elif content.type == ContentType.BULLET_LIST:
            text_box = ppt_slide.shapes.add_textbox(
                Inches(1), top, Inches(8), Inches(0.3 * len(content.value))
            )
            text_frame = text_box.text_frame

            for point in content.value:
                p = text_frame.add_paragraph()
                p.text = point
                p.level = 0  # First level bullet

            return top + Inches(0.3 * len(content.value) + 0.5)

        elif content.type == ContentType.IMAGE:
            try:
                response = requests.get(content.value['url'])
                image_data = BytesIO(response.content)
                ppt_slide.shapes.add_picture(
                    image_data, Inches(1), top, width=Inches(8)
                )
                # Add caption if exists
                if content.value.get('caption'):
                    caption_top = top + Inches(4)  # Adjust based on image height
                    text_box = ppt_slide.shapes.add_textbox(
                        Inches(1), caption_top, Inches(8), Inches(0.5)
                    )
                    text_box.text_frame.text = content.value['caption']
                    return caption_top + Inches(0.7)
                return top + Inches(4.5)
            except Exception as e:
                print(f"Error adding image: {e}")
                return top

        elif content.type == ContentType.TABLE:
            rows = len(content.value['rows']) + 1  # +1 for header
            cols = len(content.value['headers'])
            shape = ppt_slide.shapes.add_table(
                rows, cols, Inches(1), top, Inches(8), Inches(0.4 * rows)
            )
            table = shape.table

            # Add headers
            for i, header in enumerate(content.value['headers']):
                table.cell(0, i).text = header

            # Add rows
            for i, row in enumerate(content.value['rows'], 1):
                for j, cell in enumerate(row):
                    table.cell(i, j).text = str(cell)

            return top + Inches(0.4 * rows + 0.5)

        # Default spacing if content type not handled
        return top + Inches(1)